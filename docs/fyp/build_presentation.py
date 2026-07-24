from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FYP_DIR = ROOT / "docs" / "fyp"
ASSETS_DIR = FYP_DIR / "assets"
OUTPUT = FYP_DIR / "Smart_Attendance_Seminar_Presentation.pptx"

UNILORIN_LOGO = ASSETS_DIR / "unilorin_logo_nobg.png"
APP_ICON = ROOT / "mobile" / "app" / "assets" / "app_icon.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(6, 24, 57)
BLUE = RGBColor(18, 83, 227)
TEAL = RGBColor(0, 176, 150)
GREEN = RGBColor(12, 160, 103)
GOLD = RGBColor(198, 158, 83)
INK = RGBColor(29, 39, 56)
MUTED = RGBColor(91, 104, 124)
LIGHT = RGBColor(247, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(218, 226, 237)
SOFT_BLUE = RGBColor(226, 238, 255)
SOFT_TEAL = RGBColor(226, 250, 244)
SOFT_GOLD = RGBColor(255, 246, 225)
SOFT_RED = RGBColor(255, 235, 235)

FONT = "Times New Roman"


def set_run(run, size=20, color=INK, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_textbox(slide, x, y, w, h, text="", size=20, color=INK, bold=False,
                align=PP_ALIGN.LEFT, valign=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold)
    if valign:
        tf.vertical_anchor = valign
    return box


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT

    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = NAVY
    stripe.line.fill.background()


def add_footer(slide, number):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.15), Inches(12.25), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    add_textbox(
        slide,
        Inches(0.65),
        Inches(7.2),
        Inches(7.0),
        Inches(0.22),
        "Smart Attendance System using Acoustic and BLE Proximity Verification",
        size=8.5,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(12.05),
        Inches(7.18),
        Inches(0.65),
        Inches(0.24),
        f"{number}/14",
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_slide_title(slide, title, subtitle=None):
    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.45),
        Inches(10.2),
        Inches(0.48),
        title,
        size=25,
        color=NAVY,
        bold=True,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.0),
        Inches(1.35),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    if subtitle:
        add_textbox(
            slide,
            Inches(2.1),
            Inches(0.93),
            Inches(8.7),
            Inches(0.28),
            subtitle,
            size=12,
            color=MUTED,
        )


def add_card(slide, x, y, w, h, fill=CARD, line=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(1)
    return card


def add_badge(slide, x, y, text, fill, color=NAVY, w=1.4):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(w), Inches(0.34)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    badge.line.fill.background()
    add_textbox(slide, x + Inches(0.08), y + Inches(0.06), Inches(w - 0.16), Inches(0.18), text, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    return badge


def add_bullets(slide, x, y, w, h, bullets, size=17, color=INK, gap=4,
                bullet_color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(gap)
        p.text = ""
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.bullet = True
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=color)
    return box


def add_numbered_list(slide, x, y, w, h, items, size=15.2):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items, start=1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.space_after = Pt(5)
        p.text = ""
        run = p.add_run()
        run.text = f"{i}. {item}"
        set_run(run, size=size, color=INK)
    return box


def add_icon_circle(slide, x, y, label, fill=SOFT_BLUE, color=BLUE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.48), Inches(0.48))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.fill.background()
    add_textbox(slide, x, y + Inches(0.1), Inches(0.48), Inches(0.2), label, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_image_if_exists(slide, path, x, y, w=None, h=None):
    if path.exists():
        if w and h:
            slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        elif w:
            slide.shapes.add_picture(str(path), x, y, width=w)
        elif h:
            slide.shapes.add_picture(str(path), x, y, height=h)


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    wave = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(8.5), Inches(0.6), Inches(4.2), Inches(4.2))
    wave.line.color.rgb = RGBColor(46, 113, 255)
    wave.line.width = Pt(5)

    add_image_if_exists(slide, UNILORIN_LOGO, Inches(0.75), Inches(0.45), h=Inches(1.05))
    add_image_if_exists(slide, APP_ICON, Inches(11.45), Inches(0.5), h=Inches(0.95))

    add_textbox(
        slide,
        Inches(0.9),
        Inches(1.85),
        Inches(11.3),
        Inches(1.25),
        "Design and Implementation of a Smart Attendance System Using Acoustic and Bluetooth Low Energy Proximity Verification",
        size=31,
        color=RGBColor(255, 255, 255),
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.95),
        Inches(3.32),
        Inches(8.7),
        Inches(0.35),
        "Final Year Project Seminar Presentation",
        size=18,
        color=RGBColor(212, 226, 255),
        bold=True,
    )

    details = [
        "Presenter: Joshua Olugbemi Iyanuoluwa",
        "Matric Number: 21/52HP071",
        "Department of Telecommunication Science",
        "University of Ilorin, Ilorin",
    ]
    y = 4.15
    for item in details:
        add_icon_circle(slide, Inches(0.95), Inches(y), "", SOFT_TEAL, TEAL)
        add_textbox(slide, Inches(1.55), Inches(y + 0.07), Inches(8.4), Inches(0.27), item, size=15, color=RGBColor(238, 246, 255))
        y += 0.48
    add_footer(slide, 1)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Introduction", "Why attendance needs stronger classroom verification")
    add_bullets(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(6.5),
        Inches(4.6),
        [
            "Attendance tracking supports academic monitoring and classroom accountability.",
            "Manual attendance is slow, stressful in large classes, and vulnerable to proxy signing.",
            "Static digital methods such as shared QR codes can still be forwarded to absent students.",
            "Modern smartphones already provide Bluetooth, microphone, speaker, and network capabilities.",
            "The project applies telecommunication concepts to practical attendance verification.",
        ],
        size=17.2,
    )

    add_card(slide, Inches(7.75), Inches(1.55), Inches(4.75), Inches(4.5), fill=RGBColor(241, 247, 255))
    add_textbox(slide, Inches(8.1), Inches(1.9), Inches(4.1), Inches(0.45), "Presence Verification", size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, Inches(8.35), Inches(3.8), Inches(9.65), Inches(3.8), TEAL)
    add_arrow(slide, Inches(10.55), Inches(3.8), Inches(11.75), Inches(3.8), TEAL)
    for x, label, sub, fill in [
        (8.0, "Lecturer", "broadcasts session signal", SOFT_BLUE),
        (9.75, "Student", "scans proof", SOFT_TEAL),
        (11.0, "Backend", "validates", SOFT_GOLD),
    ]:
        add_card(slide, Inches(x), Inches(3.0), Inches(1.25), Inches(1.2), fill=fill)
        add_textbox(slide, Inches(x + 0.08), Inches(3.18), Inches(1.08), Inches(0.25), label, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.08), Inches(3.55), Inches(1.08), Inches(0.42), sub, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Problem Statement", "The weakness is not recording attendance; it is proving presence")

    problems = [
        ("Manual signing", "Slow and easy to manipulate."),
        ("Shared QR codes", "Can be photographed or forwarded."),
        ("GPS indoors", "Often inaccurate inside classrooms."),
        ("Single signal", "May fail due to device or environment limits."),
    ]
    x_positions = [0.75, 3.85, 6.95, 10.05]
    colors = [SOFT_RED, SOFT_GOLD, SOFT_BLUE, SOFT_TEAL]
    for (title, body), x, fill in zip(problems, x_positions, colors):
        add_card(slide, Inches(x), Inches(1.55), Inches(2.55), Inches(2.0), fill=fill)
        add_textbox(slide, Inches(x + 0.18), Inches(1.85), Inches(2.2), Inches(0.35), title, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.22), Inches(2.45), Inches(2.1), Inches(0.62), body, size=13.2, color=INK, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(1.45), Inches(4.2), Inches(10.45), Inches(1.35), fill=NAVY)
    add_textbox(
        slide,
        Inches(1.8),
        Inches(4.48),
        Inches(9.75),
        Inches(0.7),
        "Core problem: a practical attendance system is needed to verify student presence using session-specific proximity evidence and backend validation.",
        size=18,
        color=RGBColor(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Aim and Objectives", "Five objectives aligned with the current implementation")

    add_card(slide, Inches(0.75), Inches(1.35), Inches(12.0), Inches(1.0), fill=SOFT_BLUE)
    add_textbox(slide, Inches(1.0), Inches(1.55), Inches(11.5), Inches(0.45), "Aim: To design and implement a smart attendance system using acoustic beaconing and BLE proximity verification for classroom attendance management, with Wi-Fi/LAN retained as a secondary fallback.", size=15.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    objectives = [
        "Design a mobile-based smart attendance system for creating and broadcasting session-specific attendance signals.",
        "Implement proximity verification using acoustic beaconing and Bluetooth Low Energy.",
        "Develop backend validation using session identity, freshness, duplicate prevention, replay protection, and device trust.",
        "Provide student and lecturer interfaces for scanning, submission, reports, history, and CSV export.",
        "Evaluate the system on real Android devices under different signal conditions and classroom-like scenarios.",
    ]
    add_numbered_list(slide, Inches(1.0), Inches(2.75), Inches(11.25), Inches(3.7), objectives, size=15.6)
    add_footer(slide, 4)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Literature Review Summary", "What existing attendance technologies show")

    rows = [
        ("RFID/NFC", "Automated but requires tags/readers", SOFT_BLUE),
        ("QR Code", "Simple but can be shared", SOFT_GOLD),
        ("Biometrics", "Strong identity but privacy/cost concerns", SOFT_RED),
        ("BLE", "Good proximity signal but RSSI varies", SOFT_TEAL),
        ("Acoustic", "Uses speakers/mics but noise-sensitive", RGBColor(235, 245, 255)),
        ("Wi-Fi/LAN", "Useful fallback but not exact proximity", RGBColor(240, 243, 248)),
    ]
    y_positions = [1.45, 2.35, 3.25, 4.15, 5.05, 5.95]
    for (tech, note, fill), y in zip(rows, y_positions):
        add_card(slide, Inches(0.8), Inches(y), Inches(2.05), Inches(0.58), fill=fill)
        add_textbox(slide, Inches(0.95), Inches(y + 0.13), Inches(1.75), Inches(0.2), tech, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(3.1), Inches(y + 0.1), Inches(8.95), Inches(0.28), note, size=15.2, color=INK)

    add_badge(slide, Inches(8.6), Inches(1.0), "Key sources", SOFT_TEAL, GREEN, w=1.35)
    add_textbox(slide, Inches(9.95), Inches(0.98), Inches(2.7), Inches(0.34), "Puckdeevongs et al.; Ramirez et al.; Getreuer et al.; Kim et al.", size=9.8, color=MUTED)
    add_footer(slide, 5)


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Identified Research Gap", "A single signal is usually not enough")

    add_card(slide, Inches(0.85), Inches(1.45), Inches(3.3), Inches(3.5), fill=RGBColor(255, 245, 245))
    add_textbox(slide, Inches(1.1), Inches(1.75), Inches(2.8), Inches(0.45), "Existing weaknesses", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(1.12), Inches(2.35), Inches(2.75), Inches(1.8), ["Proxy attendance", "Code sharing", "Single-channel failure", "Weak replay control"], size=13.5, gap=3)

    add_arrow(slide, Inches(4.45), Inches(3.15), Inches(6.15), Inches(3.15), TEAL)
    add_textbox(slide, Inches(4.55), Inches(2.65), Inches(1.45), Inches(0.35), "Gap", size=19, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(6.55), Inches(1.45), Inches(5.85), Inches(3.5), fill=RGBColor(240, 249, 255))
    add_textbox(slide, Inches(6.85), Inches(1.75), Inches(5.2), Inches(0.45), "Needed system direction", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(6.95), Inches(2.3), Inches(4.9), Inches(1.95), ["Proximity evidence", "Fresh session signals", "Backend validation", "Duplicate prevention", "Device-trust checks"], size=13.5, gap=3)

    add_card(slide, Inches(1.25), Inches(5.45), Inches(10.85), Inches(0.85), fill=NAVY)
    add_textbox(slide, Inches(1.65), Inches(5.65), Inches(10.1), Inches(0.35), "The project addresses this gap using BLE, acoustic proof, Wi-Fi fallback, and server-side validation.", size=16.5, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Proposed System Overview", "From session broadcast to validated attendance record")

    blocks = [
        ("Lecturer App", "Create session\nStart broadcast", SOFT_BLUE),
        ("BLE + Acoustic", "Session nonce\nShort token", SOFT_TEAL),
        ("Student App", "Scan proof\nSubmit evidence", SOFT_GOLD),
        ("Django Backend", "Validate proof\nStore record", RGBColor(236, 242, 255)),
        ("Report/CSV", "Session report\nExport records", RGBColor(239, 255, 246)),
    ]
    x = 0.65
    y = 2.3
    for i, (title, body, fill) in enumerate(blocks):
        add_card(slide, Inches(x), Inches(y), Inches(2.15), Inches(1.25), fill=fill)
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.18), Inches(1.9), Inches(0.25), title, size=13.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.57), Inches(1.85), Inches(0.45), body, size=10.8, color=INK, align=PP_ALIGN.CENTER)
        if i < len(blocks) - 1:
            add_arrow(slide, Inches(x + 2.18), Inches(y + 0.62), Inches(x + 2.7), Inches(y + 0.62), TEAL)
        x += 2.55

    add_bullets(
        slide,
        Inches(1.2),
        Inches(4.45),
        Inches(10.7),
        Inches(1.25),
        [
            "BLE is the main practical classroom-range signal in the current prototype.",
            "Acoustic proof is retained as short-range copresence evidence.",
            "Wi-Fi/LAN is a fallback only for controlled classroom-network situations.",
        ],
        size=15.5,
        gap=3,
    )
    add_footer(slide, 7)


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "System Architecture", "Client-server design with native Android signal services")

    layers = [
        ("Flutter Mobile App", "Lecturer UI, student UI, scan submission, reports", SOFT_BLUE),
        ("Native Android Layer", "BLE advertising/scanning, acoustic encode/decode, runtime permissions", SOFT_TEAL),
        ("Django Backend API", "Authentication, sessions, proof validation, duplicate checks, CSV export", SOFT_GOLD),
        ("Database", "Users, devices, sessions, attendance records, validation metadata", RGBColor(239, 255, 246)),
    ]
    y = 1.5
    for title, body, fill in layers:
        add_card(slide, Inches(1.2), Inches(y), Inches(10.9), Inches(0.85), fill=fill)
        add_textbox(slide, Inches(1.55), Inches(y + 0.15), Inches(2.7), Inches(0.25), title, size=16, color=NAVY, bold=True)
        add_textbox(slide, Inches(4.45), Inches(y + 0.17), Inches(7.1), Inches(0.25), body, size=13.8, color=INK)
        y += 1.15

    add_textbox(slide, Inches(1.15), Inches(6.15), Inches(11.0), Inches(0.35), "Validation is handled by the backend so the mobile app cannot simply trust itself.", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Methodology", "Agile development with real-device feedback")

    phases = [
        "Backend\nsessions",
        "Mobile\ninterfaces",
        "Acoustic\nprototype",
        "BLE\nproximity",
        "Validation\nreports",
        "Testing and\nrefinement",
    ]
    x = 0.9
    for i, phase in enumerate(phases, start=1):
        add_card(slide, Inches(x), Inches(2.0), Inches(1.75), Inches(1.15), fill=SOFT_BLUE if i % 2 else SOFT_TEAL)
        add_textbox(slide, Inches(x + 0.18), Inches(2.18), Inches(1.38), Inches(0.52), phase, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_badge(slide, Inches(x + 0.5), Inches(3.0), f"{i}", RGBColor(255, 255, 255), BLUE, w=0.55)
        if i < len(phases):
            add_arrow(slide, Inches(x + 1.78), Inches(2.58), Inches(x + 2.28), Inches(2.58), TEAL)
        x += 2.0

    add_card(slide, Inches(1.1), Inches(4.45), Inches(11.05), Inches(1.25), fill=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(1.55), Inches(4.7), Inches(10.15), Inches(0.5), "Agile was suitable because testing changed the design: acoustic became a short-range proof, BLE became the main classroom signal, and Wi-Fi/LAN remained a fallback.", size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Attendance Verification Workflow", "How a student's attendance proof is accepted")

    items = [
        "Lecturer creates a session.",
        "Lecturer starts BLE/acoustic broadcast.",
        "Student logs in on a trusted device.",
        "Student scans for BLE or acoustic proof.",
        "Wi-Fi/LAN fallback may be used where allowed.",
        "Backend validates session, freshness, replay status, duplicates, and device trust.",
        "Attendance is accepted, rejected, or flagged.",
    ]
    add_numbered_list(slide, Inches(0.95), Inches(1.45), Inches(5.8), Inches(4.95), items, size=15)

    add_card(slide, Inches(7.25), Inches(1.55), Inches(4.8), Inches(4.55), fill=RGBColor(241, 247, 255))
    y = 2.0
    checks = [("Session", "active"), ("Freshness", "within time"), ("Replay", "not reused"), ("Duplicate", "not submitted"), ("Device", "trusted/checked")]
    for label, status in checks:
        add_icon_circle(slide, Inches(7.65), Inches(y), "OK", SOFT_TEAL, GREEN)
        add_textbox(slide, Inches(8.28), Inches(y + 0.08), Inches(1.65), Inches(0.25), label, size=14.5, color=NAVY, bold=True)
        add_textbox(slide, Inches(9.7), Inches(y + 0.08), Inches(1.8), Inches(0.25), status, size=13, color=MUTED)
        y += 0.68
    add_footer(slide, 10)


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Key Implementation Features", "Main features currently available in the prototype")

    features = [
        ("Sessions", "Create and manage live attendance sessions."),
        ("BLE", "Advertise and scan session-specific proximity nonce."),
        ("Acoustic", "Transmit and decode short attendance token."),
        ("Wi-Fi/LAN", "Fallback proof for controlled networks."),
        ("Anti-fraud", "Device binding, duplicate prevention, replay checks."),
        ("Reports", "Validation report and CSV export by session."),
    ]
    x_positions = [0.85, 4.65, 8.45]
    y_positions = [1.55, 3.65]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            title, body = features[idx]
            add_card(slide, Inches(x), Inches(y), Inches(3.1), Inches(1.45), fill=RGBColor(255, 255, 255))
            add_badge(slide, Inches(x + 0.22), Inches(y + 0.18), title, SOFT_BLUE if idx % 2 == 0 else SOFT_TEAL, NAVY, w=1.15)
            add_textbox(slide, Inches(x + 0.25), Inches(y + 0.7), Inches(2.55), Inches(0.45), body, size=12.7, color=INK, align=PP_ALIGN.CENTER)
            idx += 1
    add_footer(slide, 11)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Preliminary Results and Observations", "What has been confirmed from real-device testing")

    add_card(slide, Inches(0.75), Inches(1.35), Inches(5.65), Inches(4.7), fill=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(1.05), Inches(1.6), Inches(5.0), Inches(0.35), "Current working results", size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.05), Inches(2.05), Inches(5.05), Inches(2.9), [
        "Android app runs on real phones.",
        "Lecturer can create sessions and broadcast proof signals.",
        "Student can scan and submit attendance proof.",
        "BLE range improved after Location, Nearby Devices, and Bluetooth permissions were enabled.",
        "Reports and CSV export are available.",
    ], size=13.8, gap=2)

    add_card(slide, Inches(6.85), Inches(1.35), Inches(5.65), Inches(4.7), fill=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(7.15), Inches(1.6), Inches(5.0), Inches(0.35), "Signal role summary", size=18, color=NAVY, bold=True)
    table_data = [
        ("BLE", "Main practical classroom signal"),
        ("Acoustic", "Short-range copresence proof"),
        ("Wi-Fi/LAN", "Fallback support only"),
        ("BLE beacon", "Future larger-room extension"),
    ]
    y = 2.1
    for name, role in table_data:
        add_card(slide, Inches(7.15), Inches(y), Inches(1.35), Inches(0.5), fill=SOFT_TEAL if name == "BLE" else SOFT_BLUE)
        add_textbox(slide, Inches(7.25), Inches(y + 0.11), Inches(1.15), Inches(0.2), name, size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(8.75), Inches(y + 0.1), Inches(3.2), Inches(0.22), role, size=12.5, color=INK)
        y += 0.72

    add_textbox(slide, Inches(1.05), Inches(6.35), Inches(11.15), Inches(0.35), "Honest position: BLE is currently stronger; acoustic works but remains short-ranged and noise-sensitive.", size=15.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 12)


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "Conclusion and Future Work", "Where the project stands and what comes next")

    add_card(slide, Inches(0.8), Inches(1.45), Inches(5.75), Inches(4.65), fill=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(1.12), Inches(1.75), Inches(5.1), Inches(0.35), "Conclusion", size=19, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.1), Inches(2.2), Inches(5.1), Inches(2.95), [
        "The system demonstrates mobile attendance verification using BLE and acoustic proximity signals.",
        "BLE is currently the strongest classroom-range proof in the prototype.",
        "Backend validation improves trust through freshness, replay protection, duplicate checks, and device trust.",
        "Lecturers can view reports and export attendance records.",
    ], size=13.5, gap=2)

    add_card(slide, Inches(6.85), Inches(1.45), Inches(5.65), Inches(4.65), fill=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(7.18), Inches(1.75), Inches(5.0), Inches(0.35), "Future work", size=19, color=NAVY, bold=True)
    add_bullets(slide, Inches(7.18), Inches(2.2), Inches(4.9), Inches(2.95), [
        "Test fixed BLE beacons in medium and large classrooms.",
        "Improve acoustic robustness and range.",
        "Run structured field tests across room size and noise levels.",
        "Strengthen device-trust rules and exception handling.",
        "Prepare for institutional pilot deployment.",
    ], size=13.5, gap=2)
    add_footer(slide, 13)


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_title(slide, "References", "Selected APA 7 references used in the seminar")

    refs = [
        "Android Developers. (n.d.). Bluetooth permissions. https://developer.android.com/develop/connectivity/bluetooth/bt-permissions",
        "Getreuer, P., Gnegy, C., Lyon, R. F., & Saurous, R. A. (2018). Ultrasonic communication using consumer hardware. IEEE Transactions on Multimedia, 20(6), 1277-1290. https://doi.org/10.1109/TMM.2017.2766049",
        "Jia, N., Shu, H., Wang, X., Xu, B., Xi, Y., Xue, C., Liu, Y., & Wang, Z. (2022). Smartphone-based social distance detection technology with near-ultrasonic signal. Sensors, 22(19), Article 7345. https://doi.org/10.3390/s22197345",
        "Kim, M., Lee, J., & Paek, J. (2018). Neutralizing BLE beacon-based electronic attendance system using signal imitation attack. IEEE Access, 6, 77921-77930. https://doi.org/10.1109/ACCESS.2018.2884488",
        "Puckdeevongs, A., Tripathi, N. K., Witayangkurn, A., & Saengudomlert, P. (2020). Classroom attendance systems based on Bluetooth Low Energy indoor positioning technology for smart campus. Information, 11(6), Article 329. https://doi.org/10.3390/info11060329",
        "Ramirez, R., Huang, C.-Y., Liao, C.-A., Lin, P.-T., Lin, H.-W., & Liang, S.-H. (2021). A practice of BLE RSSI measurement for indoor positioning. Sensors, 21(15), Article 5181. https://doi.org/10.3390/s21155181",
    ]

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.9), Inches(5.35))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.text = ""
        run = p.add_run()
        run.text = ref
        set_run(run, size=10.2, color=INK)
    add_footer(slide, 14)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder in [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
    ]:
        builder(prs)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build()
