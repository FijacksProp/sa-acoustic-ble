from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FYP_DIR = ROOT / "docs" / "fyp"
ASSETS_DIR = FYP_DIR / "assets"
OUTPUT = FYP_DIR / "Smart_Attendance_Professional_Times_New_Roman_Image_Placeholders.pptx"

UNILORIN_LOGO = ASSETS_DIR / "unilorin_logo_nobg.png"
APP_ICON = ROOT / "mobile" / "app" / "assets" / "app_icon.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Times New Roman"

NAVY = RGBColor(8, 27, 63)
BLUE = RGBColor(20, 86, 230)
TEAL = RGBColor(0, 174, 151)
GREEN = RGBColor(16, 145, 96)
GOLD = RGBColor(197, 153, 70)
INK = RGBColor(27, 39, 56)
MUTED = RGBColor(86, 100, 121)
LIGHT = RGBColor(247, 250, 253)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(217, 226, 238)
SOFT_BLUE = RGBColor(226, 238, 255)
SOFT_TEAL = RGBColor(225, 249, 244)
SOFT_GOLD = RGBColor(255, 247, 226)
SOFT_RED = RGBColor(255, 235, 235)
SOFT_GREEN = RGBColor(232, 250, 241)


def set_run(run, size=18, color=INK, bold=False, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, color=color, bold=bold, italic=italic)
    return box


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.18), SLIDE_W, Inches(0.035))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def add_card(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(0.8)
    return card


def add_footer(slide, n):
    add_text(
        slide,
        Inches(0.55),
        Inches(7.13),
        Inches(9.4),
        Inches(0.22),
        "Smart Attendance System | Dept. of Telecommunication Science | University of Ilorin",
        size=8.7,
        color=MUTED,
    )
    add_text(
        slide,
        Inches(11.85),
        Inches(7.13),
        Inches(0.95),
        Inches(0.22),
        f"{n}/11",
        size=8.7,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.65), Inches(0.42), Inches(11.2), Inches(0.45), title, size=24, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, Inches(0.68), Inches(0.88), Inches(11.6), Inches(0.28), subtitle, size=11.2, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(1.18), Inches(1.45), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_bullets(slide, x, y, w, h, items, size=15.5, color=INK, gap=4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(gap)
        p.bullet = True
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=color)
    return box


def add_numbered(slide, x, y, w, h, items, size=14.5, gap=4):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items, start=1):
        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.text = ""
        run = p.add_run()
        run.text = f"{idx}. {item}"
        set_run(run, size=size, color=INK)
    return box


def add_badge(slide, x, y, text, fill=SOFT_BLUE, color=NAVY, w=1.25):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(w), Inches(0.33))
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    badge.line.fill.background()
    add_text(slide, x + Inches(0.08), y + Inches(0.065), Inches(w - 0.16), Inches(0.16), text, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def add_image(slide, path, x, y, w=None, h=None):
    if path.exists():
        if w and h:
            slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        elif w:
            slide.shapes.add_picture(str(path), x, y, width=w)
        elif h:
            slide.shapes.add_picture(str(path), x, y, height=h)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_image(slide, UNILORIN_LOGO, Inches(0.7), Inches(0.55), h=Inches(1.05))
    add_image(slide, APP_ICON, Inches(11.5), Inches(0.68), h=Inches(0.85))

    add_text(slide, Inches(1.7), Inches(0.62), Inches(9.8), Inches(0.45), "UNIVERSITY OF ILORIN, ILORIN", size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.15), Inches(1.08), Inches(8.9), Inches(0.3), "Faculty of Communication and Information Sciences", size=12.5, color=RGBColor(220, 232, 255), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.15), Inches(1.42), Inches(8.9), Inches(0.3), "Department of Telecommunication Science", size=12.5, color=RGBColor(220, 232, 255), align=PP_ALIGN.CENTER)

    add_badge(slide, Inches(5.35), Inches(2.05), "PROJECT PRESENTATION", SOFT_TEAL, NAVY, w=2.7)
    add_text(
        slide,
        Inches(1.1),
        Inches(2.75),
        Inches(11.1),
        Inches(1.15),
        "Design and Implementation of a Smart Attendance System Using Acoustic and Bluetooth Low Energy Proximity Verification",
        size=29,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, Inches(2.2), Inches(4.42), Inches(8.9), Inches(0.33), "Presented by:", size=13.5, color=RGBColor(220, 232, 255), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.2), Inches(4.78), Inches(8.9), Inches(0.34), "OLUGBEMI Joshua Iyanuoluwa - 21/52HP071", size=16.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.2), Inches(5.38), Inches(8.9), Inches(0.3), "Project Supervisor: Dr. IMAM FULANI", size=14, color=RGBColor(230, 239, 255), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.2), Inches(6.02), Inches(8.9), Inches(0.28), "May, 2026", size=13.2, color=RGBColor(230, 239, 255), align=PP_ALIGN.CENTER)


def slide_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "INTRODUCTION", "Background of the study")

    add_card(slide, Inches(0.75), Inches(1.45), Inches(5.8), Inches(4.95), fill=WHITE)
    add_text(slide, Inches(1.05), Inches(1.72), Inches(5.15), Inches(0.35), "Attendance remains operationally difficult", size=19, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.05), Inches(2.22), Inches(5.1), Inches(3.35), [
        "Manual roll call consumes class time and creates avoidable administrative delay.",
        "Paper registers provide weak tamper-evidence and can be manipulated.",
        "Static QR codes and PINs can be shared outside the classroom.",
        "Most basic systems verify user identity, but not physical presence.",
    ], size=15.2, gap=4)

    add_card(slide, Inches(6.95), Inches(1.45), Inches(5.65), Inches(4.95), fill=RGBColor(240, 247, 255))
    add_text(slide, Inches(7.25), Inches(1.72), Inches(5.05), Inches(0.35), "Telecommunication perspective", size=19, color=NAVY, bold=True)
    add_bullets(slide, Inches(7.25), Inches(2.22), Inches(5.05), Inches(3.35), [
        "The lecturer acts as a signal broadcaster for an active attendance session.",
        "The student device acts as a receiver that captures a short-lived proof signal.",
        "BLE provides the main practical classroom-range proof.",
        "Acoustic beaconing provides short-range copresence evidence.",
        "The backend validates freshness, replay status, duplicate status, and device trust.",
    ], size=14.6, gap=3)
    add_footer(slide, 2)


def slide_lit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "LITERATURE REVIEW", "Verified related works and the gap they leave")

    table = slide.shapes.add_table(6, 5, Inches(0.35), Inches(1.45), Inches(12.65), Inches(5.28)).table
    widths = [0.45, 2.35, 2.05, 3.6, 4.2]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    headers = ["S/N", "Source", "Area", "Work Done", "Weakness / Gap"]
    rows = [
        ["1", "Ayop et al. (2018)", "QR + GPS", "Developed location-aware event attendance using QR code and GPS.", "QR can be shared; GPS is less reliable indoors."],
        ["2", "Getreuer et al. (2018)", "Acoustic", "Demonstrated ultrasonic communication using consumer hardware.", "Not designed for attendance; range depends on speaker/mic and noise."],
        ["3", "Kim et al. (2018)", "BLE security", "Analyzed BLE beacon attendance and signal imitation attacks.", "BLE proof needs freshness and replay protection."],
        ["4", "Puckdeevongs et al. (2020)", "BLE attendance", "Proposed BLE indoor-positioning attendance for smart campus.", "RSSI and beacon coverage vary across rooms/devices."],
        ["5", "Ramirez et al. (2021)", "BLE RSSI", "Studied practical BLE RSSI measurement for indoor positioning.", "RSSI supports proximity, not exact classroom distance."],
    ]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=10.4, color=WHITE, bold=True)
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(243, 248, 255)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                for run in p.runs:
                    set_run(run, size=9.4 if c > 0 else 10.2, color=INK, bold=(c == 0))
    add_text(slide, Inches(0.7), Inches(6.87), Inches(12.0), Inches(0.22), "Summary: existing works support automation and proximity verification, but single-channel systems still require freshness, replay protection, and backend validation.", size=10.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "PROBLEM STATEMENT", "Why the project is necessary")
    add_card(slide, Inches(0.85), Inches(1.55), Inches(11.75), Inches(2.2), fill=WHITE)
    problem = (
        "Attendance in many higher-education settings still relies on manual roll calls, paper sheets, "
        "or basic digital tools such as static QR codes and PINs. These methods waste class time, offer "
        "little protection against tampering, and make it easy to copy or reuse attendance artifacts outside "
        "the classroom."
    )
    add_text(slide, Inches(1.25), Inches(1.95), Inches(10.95), Inches(1.08), problem, size=18, color=INK, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(0.85), Inches(4.05), Inches(11.75), Inches(1.55), fill=NAVY)
    problem2 = (
        "The core problem is that most systems verify login identity, but not physical presence, signal freshness, "
        "or replay resistance. A more secure approach needs time-limited proximity proof and server-side validation."
    )
    add_text(slide, Inches(1.25), Inches(4.38), Inches(10.95), Inches(0.72), problem2, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)


def slide_aim(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "AIM AND OBJECTIVES", "Updated to match the implemented prototype")
    add_card(slide, Inches(0.75), Inches(1.38), Inches(12.0), Inches(0.95), fill=SOFT_BLUE)
    add_text(slide, Inches(1.02), Inches(1.62), Inches(11.45), Inches(0.36), "Aim: To design and implement a secure, role-aware smart attendance system that reduces fraud and operational delay by validating signal-based proof of physical presence.", size=15.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    objectives = [
        "Design a mobile-based smart attendance system for creating and broadcasting session-specific attendance signals.",
        "Implement proximity verification using acoustic beaconing and Bluetooth Low Energy, with BLE as the main classroom-range signal and acoustic proof as short-range copresence evidence.",
        "Develop backend validation using session identity, signal freshness, duplicate prevention, replay protection, and device-trust checks.",
        "Provide student and lecturer interfaces for scanning, proof submission, attendance history, validation reports, and CSV export.",
        "Evaluate the system on real Android devices under different signal conditions, distances, permission states, and classroom-like scenarios.",
    ]
    add_numbered(slide, Inches(1.0), Inches(2.72), Inches(11.3), Inches(3.7), objectives, size=15.0, gap=4)
    add_footer(slide, 5)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "PROPOSED METHODOLOGY: SYSTEM ARCHITECTURE", "Current implemented architecture")

    blocks = [
        ("CLIENT TIER - Flutter Android App", "Lecturer dashboard, student scan screen, profile, reports, backend URL setting", SOFT_BLUE),
        ("NATIVE ANDROID SIGNAL SERVICES", "BLE advertising/scanning, acoustic encoding/decoding, microphone/location/Bluetooth permissions", SOFT_TEAL),
        ("BACKEND TIER - Django REST API", "Authentication, session control, proof submission, validation, CSV export", SOFT_GOLD),
        ("DATA TIER - SQLite/PostgreSQL-ready", "Users, devices, sessions, attendance proofs, replay guard, validation metadata", SOFT_GREEN),
    ]
    y = 1.48
    for title, desc, fill in blocks:
        add_card(slide, Inches(1.0), Inches(y), Inches(11.35), Inches(0.85), fill=fill)
        add_text(slide, Inches(1.35), Inches(y + 0.12), Inches(3.9), Inches(0.25), title, size=14.6, color=NAVY, bold=True)
        add_text(slide, Inches(5.4), Inches(y + 0.14), Inches(6.6), Inches(0.25), desc, size=12.7, color=INK)
        y += 1.06
    add_arrow(slide, Inches(6.65), Inches(2.35), Inches(6.65), Inches(2.55))
    add_arrow(slide, Inches(6.65), Inches(3.41), Inches(6.65), Inches(3.61))
    add_arrow(slide, Inches(6.65), Inches(4.47), Inches(6.65), Inches(4.67))
    add_text(slide, Inches(1.15), Inches(6.22), Inches(11.1), Inches(0.35), "Validation is intentionally server-side: the app captures proof, while the backend decides whether the attendance record is valid.", size=15.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def slide_identity_payload(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "PROPOSED METHODOLOGY: IDENTITY, PAYLOAD & PROOF CONSTRUCTION", "Rewritten to match the current implementation")

    add_card(slide, Inches(0.55), Inches(1.38), Inches(3.9), Inches(4.8), fill=WHITE)
    add_text(slide, Inches(0.85), Inches(1.62), Inches(3.3), Inches(0.3), "Identity & Access Control", size=16.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(0.85), Inches(2.08), Inches(3.25), Inches(3.35), [
        "Students authenticate with matric number and password.",
        "Lecturers authenticate with lecturer credentials.",
        "API endpoints are protected by token-based authentication.",
        "Lecturers view only their sessions and validation reports.",
        "Device ID binding is used to reduce account-sharing abuse.",
    ], size=12.4, gap=2)

    add_card(slide, Inches(4.72), Inches(1.38), Inches(3.9), Inches(4.8), fill=RGBColor(240, 247, 255))
    add_text(slide, Inches(5.02), Inches(1.62), Inches(3.3), Inches(0.3), "Signal Payload Model", size=16.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(5.02), Inches(2.08), Inches(3.25), Inches(3.35), [
        "BLE proof contains a session-specific short-lived nonce.",
        "Acoustic proof contains a short session token broadcast through sound.",
        "Wi-Fi/LAN proof is retained only as a controlled fallback.",
        "Signals are checked against active session identity and freshness.",
        "BLE is treated as the main practical classroom signal.",
    ], size=12.4, gap=2)

    add_card(slide, Inches(8.88), Inches(1.38), Inches(3.9), Inches(4.8), fill=RGBColor(239, 255, 246))
    add_text(slide, Inches(9.18), Inches(1.62), Inches(3.3), Inches(0.3), "Proof Construction", size=16.2, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(9.18), Inches(2.08), Inches(3.25), Inches(3.35), [
        "Student ID comes from the authenticated session.",
        "Device ID comes from the student's registered device profile.",
        "Scan mode is recorded as BLE, acoustic, or Wi-Fi/LAN fallback.",
        "Submitted proof is validated before an attendance record is stored.",
        "Duplicate and replay attempts are rejected or flagged.",
    ], size=12.4, gap=2)
    add_footer(slide, 7)


def slide_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "WORKFLOW", "Attendance verification workflow")

    steps = [
        ("1", "Lecturer creates an active attendance session."),
        ("2", "Lecturer starts BLE and acoustic broadcast."),
        ("3", "Student logs in on a registered/trusted device."),
        ("4", "Student scans BLE or acoustic proof in the classroom."),
        ("5", "Wi-Fi/LAN fallback may be used only under controlled conditions."),
        ("6", "App submits proof to the backend API."),
        ("7", "Backend validates session, freshness, duplicates, replay status, and device trust."),
        ("8", "Attendance is accepted, rejected, or flagged."),
    ]
    y = 1.43
    for idx, (num, text) in enumerate(steps):
        x = 0.75 if idx < 4 else 6.95
        row_y = y + (idx % 4) * 1.18
        add_card(slide, Inches(x), Inches(row_y), Inches(5.55), Inches(0.85), fill=WHITE if idx % 2 == 0 else RGBColor(241, 247, 255))
        add_badge(slide, Inches(x + 0.18), Inches(row_y + 0.25), num, SOFT_TEAL, NAVY, w=0.45)
        add_text(slide, Inches(x + 0.78), Inches(row_y + 0.18), Inches(4.45), Inches(0.28), text, size=13.3, color=INK)
    add_text(slide, Inches(0.95), Inches(6.38), Inches(11.45), Inches(0.32), "Validation rules turn raw signal detection into attendance decisions.", size=15.0, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)


def slide_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "PROPOSED METHODOLOGY: BACKEND VALIDATION & SECURITY CONTROLS", "Sequential checks before a proof is recorded")

    controls = [
        ("1", "Signal freshness", "Reject stale BLE/acoustic/Wi-Fi proof outside the allowed session window."),
        ("2", "Format correctness", "Reject malformed, truncated, or unsupported proof values before recording."),
        ("3", "Session consistency", "Decoded session identity must match the active session being submitted."),
        ("4", "Session state", "Only active lecturer-owned sessions can accept attendance."),
        ("5", "Duplicate prevention", "One accepted attendance record per student per session."),
        ("6", "Replay protection", "Previously used proof values are rejected to prevent token reuse."),
        ("7", "Device trust", "Registered device ID is checked to reduce account-sharing fraud."),
    ]
    x_positions = [0.55, 4.65, 8.75]
    y_positions = [1.45, 3.35, 5.25]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            if idx >= len(controls):
                break
            num, title, desc = controls[idx]
            add_card(slide, Inches(x), Inches(y), Inches(3.55), Inches(1.35), fill=WHITE if idx % 2 == 0 else RGBColor(242, 248, 255))
            add_badge(slide, Inches(x + 0.18), Inches(y + 0.17), num, SOFT_TEAL, NAVY, w=0.45)
            add_text(slide, Inches(x + 0.78), Inches(y + 0.18), Inches(2.45), Inches(0.24), title, size=13.6, color=NAVY, bold=True)
            add_text(slide, Inches(x + 0.22), Inches(y + 0.62), Inches(3.05), Inches(0.42), desc, size=10.2, color=INK, align=PP_ALIGN.CENTER)
            idx += 1
    add_text(slide, Inches(4.7), Inches(5.78), Inches(7.55), Inches(0.4), "Security principle: physical signal capture alone is not enough; the backend must validate identity, time, session, duplicate status, replay status, and device trust.", size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "CONCLUSION", "Summary of implemented system and next steps")

    add_card(slide, Inches(0.75), Inches(1.45), Inches(5.75), Inches(4.6), fill=WHITE)
    add_text(slide, Inches(1.05), Inches(1.72), Inches(5.1), Inches(0.3), "Conclusion", size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.05), Inches(2.16), Inches(5.1), Inches(2.95), [
        "The project demonstrates mobile attendance using BLE and acoustic proximity verification.",
        "BLE is currently the strongest practical signal for classroom-range attendance proof.",
        "Acoustic proof adds short-range copresence evidence.",
        "Backend validation improves trust through freshness, duplicate prevention, replay protection, and device checks.",
        "Lecturer reports and CSV export support practical academic use.",
    ], size=12.8, gap=2)

    add_card(slide, Inches(6.85), Inches(1.45), Inches(5.75), Inches(4.6), fill=RGBColor(239, 255, 246))
    add_text(slide, Inches(7.15), Inches(1.72), Inches(5.1), Inches(0.3), "Future work", size=18, color=NAVY, bold=True)
    add_bullets(slide, Inches(7.15), Inches(2.16), Inches(5.05), Inches(2.95), [
        "Test fixed BLE beacons for medium and large classrooms.",
        "Improve acoustic signal robustness and distance.",
        "Conduct structured field testing across room size, noise, and phone models.",
        "Refine device-trust and exception scoring rules.",
        "Prepare hosted-backend deployment for easier demonstrations.",
    ], size=12.8, gap=2)
    add_footer(slide, 10)


def slide_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "REFERENCES", "Selected real sources supporting the project")

    refs_left = [
        "Android Developers. (n.d.). Bluetooth permissions. https://developer.android.com/develop/connectivity/bluetooth/bt-permissions",
        "Ayop, Z., Lin, C. Y., Anawar, S., Hamid, E., & Azhar, M. S. (2018). Location-aware event attendance system using QR code and GPS technology. International Journal of Advanced Computer Science and Applications, 9(9), 466-473. https://doi.org/10.14569/IJACSA.2018.090959",
        "Getreuer, P., Gnegy, C., Lyon, R. F., & Saurous, R. A. (2018). Ultrasonic communication using consumer hardware. IEEE Transactions on Multimedia, 20(6), 1277-1290. https://doi.org/10.1109/TMM.2017.2766049",
        "Kim, M., Lee, J., & Paek, J. (2018). Neutralizing BLE beacon-based electronic attendance system using signal imitation attack. IEEE Access, 6, 77921-77930. https://doi.org/10.1109/ACCESS.2018.2884488",
    ]
    refs_right = [
        "Nwabuwe, A., Sanghera, B., Alade, T., & Olajide, F. (2023). Fraud mitigation in attendance monitoring systems using dynamic QR code, geofencing and IMEI technologies. International Journal of Advanced Computer Science and Applications, 14(4), 938-945. https://doi.org/10.14569/IJACSA.2023.01404104",
        "Puckdeevongs, A., Tripathi, N. K., Witayangkurn, A., & Saengudomlert, P. (2020). Classroom attendance systems based on Bluetooth Low Energy indoor positioning technology for smart campus. Information, 11(6), Article 329. https://doi.org/10.3390/info11060329",
        "Ramirez, R., Huang, C.-Y., Liao, C.-A., Lin, P.-T., Lin, H.-W., & Liang, S.-H. (2021). A practice of BLE RSSI measurement for indoor positioning. Sensors, 21(15), Article 5181. https://doi.org/10.3390/s21155181",
        "Jia, N., Shu, H., Wang, X., Xu, B., Xi, Y., Xue, C., Liu, Y., & Wang, Z. (2022). Smartphone-based social distance detection technology with near-ultrasonic signal. Sensors, 22(19), Article 7345. https://doi.org/10.3390/s22197345",
    ]
    add_card(slide, Inches(0.55), Inches(1.42), Inches(6.0), Inches(5.35), fill=WHITE)
    add_card(slide, Inches(6.78), Inches(1.42), Inches(6.0), Inches(5.35), fill=WHITE)

    def add_refs(refs, x):
        y = 1.7
        for ref in refs:
            add_text(slide, Inches(x), Inches(y), Inches(5.35), Inches(0.78), ref, size=8.0, color=INK)
            y += 1.2

    add_refs(refs_left, 0.85)
    add_refs(refs_right, 7.08)
    add_footer(slide, 11)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder in [
        slide_title,
        slide_intro,
        slide_lit,
        slide_problem,
        slide_aim,
        slide_architecture,
        slide_identity_payload,
        slide_workflow,
        slide_validation,
        slide_conclusion,
        slide_references,
    ]:
        builder(prs)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


if __name__ == "__main__":
    build()
